import os
import platform
import uuid
from typing import Dict, List, Optional, Any, Union
from abc import ABC, abstractmethod


class PowerPointAdapter(ABC):
    """Abstract base class for PowerPoint automation adapters"""

    @abstractmethod
    def initialize(self) -> bool:
        """Initialize connection to PowerPoint"""
        pass

    @abstractmethod
    def get_open_presentations(self) -> List[Dict[str, Any]]:
        """Get all currently open presentations"""
        pass

    @abstractmethod
    def open_presentation(self, path: str) -> Dict[str, Any]:
        """Open a presentation from file"""
        pass

    @abstractmethod
    def create_presentation(self) -> Dict[str, Any]:
        """Create a new presentation"""
        pass

    @abstractmethod
    def save_presentation(
        self, presentation_id: str, path: str = None
    ) -> Dict[str, Any]:
        """Save a presentation"""
        pass

    @abstractmethod
    def close_presentation(
        self, presentation_id: str, save: bool = True
    ) -> Dict[str, Any]:
        """Close a presentation"""
        pass

    @abstractmethod
    def get_slides(self, presentation_id: str) -> List[Dict[str, Any]]:
        """Get slides in a presentation"""
        pass

    @abstractmethod
    def add_slide(self, presentation_id: str, layout_type: int = 1) -> Dict[str, Any]:
        """Add a new slide"""
        pass

    @abstractmethod
    def get_slide_text(self, presentation_id: str, slide_id: int) -> Dict[str, Any]:
        """Get text content from a slide"""
        pass

    @abstractmethod
    def update_text(
        self, presentation_id: str, slide_id: str, shape_id: str, text: str
    ) -> Dict[str, Any]:
        """Update text in a shape"""
        pass

    @abstractmethod
    def add_text_box(
        self,
        presentation_id: str,
        slide_id: str,
        text: str,
        left: float = 100,
        top: float = 100,
        width: float = 400,
        height: float = 200,
    ) -> Dict[str, Any]:
        """Add a text box to a slide"""
        pass

    @abstractmethod
    def set_slide_title(
        self, presentation_id: str, slide_id: str, title: str
    ) -> Dict[str, Any]:
        """Set the title of a slide"""
        pass


class WindowsCOMAdapter(PowerPointAdapter):
    """Windows COM API adapter for direct PowerPoint automation"""

    def __init__(self):
        self.ppt_app = None
        self.presentations = {}

        # Import Windows-specific modules
        try:
            import win32com.client

            self.win32com = win32com
            self.available = True
        except ImportError:
            self.available = False

    def initialize(self) -> bool:
        if not self.available:
            return False

        try:
            # Try to connect to a running PowerPoint instance
            self.ppt_app = self.win32com.client.GetActiveObject(
                "PowerPoint.Application"
            )
            return True
        except:
            try:
                # If no instance is running, create a new one
                self.ppt_app = self.win32com.client.Dispatch("PowerPoint.Application")
                self.ppt_app.Visible = True
                return True
            except:
                return False

    def get_open_presentations(self) -> List[Dict[str, Any]]:
        result = []
        if not self.ppt_app:
            self.initialize()

        if self.ppt_app:
            for i in range(1, self.ppt_app.Presentations.Count + 1):
                pres = self.ppt_app.Presentations.Item(i)
                pres_id = str(uuid.uuid4())
                self.presentations[pres_id] = pres
                result.append(
                    {
                        "id": pres_id,
                        "name": (
                            os.path.basename(pres.FullName)
                            if pres.FullName
                            else "Untitled"
                        ),
                        "path": pres.FullName,
                        "slide_count": pres.Slides.Count,
                    }
                )
        return result

    def open_presentation(self, path: str) -> Dict[str, Any]:
        if not self.ppt_app:
            self.initialize()

        if not os.path.exists(path):
            return {"error": f"File not found: {path}"}

        try:
            pres = self.ppt_app.Presentations.Open(path)
            pres_id = str(uuid.uuid4())
            self.presentations[pres_id] = pres

            return {
                "id": pres_id,
                "name": os.path.basename(path),
                "path": path,
                "slide_count": pres.Slides.Count,
            }
        except Exception as e:
            return {"error": str(e)}

    def create_presentation(self) -> Dict[str, Any]:
        if not self.ppt_app:
            self.initialize()

        try:
            pres = self.ppt_app.Presentations.Add()
            pres_id = str(uuid.uuid4())
            self.presentations[pres_id] = pres

            return {
                "id": pres_id,
                "name": "New Presentation",
                "path": "",
                "slide_count": pres.Slides.Count,
            }
        except Exception as e:
            return {"error": str(e)}

    def save_presentation(
        self, presentation_id: str, path: str = None
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]

        try:
            if path:
                pres.SaveAs(path)
            else:
                pres.Save()
            return {"success": True, "path": path if path else pres.FullName}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def close_presentation(
        self, presentation_id: str, save: bool = True
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]

        try:
            pres.Close(save)
            del self.presentations[presentation_id]
            return {"success": True}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_slides(self, presentation_id: str) -> List[Dict[str, Any]]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]
        slides = []

        try:
            slide_count = pres.Slides.Count

            for i in range(1, slide_count + 1):
                slide = pres.Slides.Item(i)
                slide_id = str(i)

                slides.append(
                    {
                        "id": slide_id,
                        "index": i,
                        "title": self._get_slide_title(slide),
                        "shape_count": slide.Shapes.Count,
                    }
                )

            return slides
        except Exception as e:
            return {"error": f"Error getting slides: {str(e)}"}

    def _get_slide_title(self, slide):
        """Helper function to extract slide title"""
        try:
            # Check for title placeholder
            for shape in slide.Shapes:
                if shape.Type == 14:  # msoPlaceholder
                    if shape.PlaceholderFormat.Type == 1:  # ppPlaceholderTitle
                        if hasattr(shape, "TextFrame") and hasattr(
                            shape.TextFrame, "TextRange"
                        ):
                            return shape.TextFrame.TextRange.Text

            # Check any text shape
            for shape in slide.Shapes:
                if (
                    shape.Type == 17
                    and hasattr(shape, "TextFrame")
                    and hasattr(shape.TextFrame, "TextRange")
                ):
                    try:
                        text = shape.TextFrame.TextRange.Text
                        if text and text.strip():
                            return text
                    except:
                        continue
        except:
            pass

        return "Untitled Slide"

    def add_slide(self, presentation_id: str, layout_type: int = 1) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]

        try:
            slide_index = pres.Slides.Count + 1
            slide = pres.Slides.Add(slide_index, layout_type)

            return {
                "id": str(slide_index),
                "index": slide_index,
                "title": "New Slide",
                "shape_count": slide.Shapes.Count,
            }
        except Exception as e:
            return {"error": f"Error adding slide: {str(e)}"}

    def get_slide_text(self, presentation_id: str, slide_id: int) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": f"Presentation ID not found: {presentation_id}"}

        pres = self.presentations[presentation_id]

        try:
            slide_count = pres.Slides.Count
        except Exception as e:
            return {"error": f"Unable to get slide count: {str(e)}"}

        if slide_count == 0:
            return {"error": "Presentation has no slides"}

        if slide_id < 1 or slide_id > slide_count:
            return {
                "error": f"Invalid slide ID: {slide_id}. Valid range is 1-{slide_count}"
            }

        try:
            slide = pres.Slides.Item(int(slide_id))
        except Exception as e:
            return {"error": f"Error retrieving slide: {str(e)}"}

        text_content = {}

        try:
            shape_count = slide.Shapes.Count
        except Exception as e:
            return {"error": f"Unable to get shape count: {str(e)}"}

        for shape_idx in range(1, shape_count + 1):
            try:
                shape = slide.Shapes.Item(shape_idx)
                shape_id = str(shape_idx)

                has_text = False
                text = ""

                try:
                    if hasattr(shape, "TextFrame2") and shape.TextFrame2.HasText:
                        has_text = True
                        text = shape.TextFrame2.TextRange.Text
                    elif (
                        hasattr(shape, "TextFrame")
                        and hasattr(shape.TextFrame, "HasText")
                        and shape.TextFrame.HasText
                    ):
                        has_text = True
                        text = shape.TextFrame.TextRange.Text
                    elif hasattr(shape, "TextFrame") and hasattr(
                        shape.TextFrame, "TextRange"
                    ):
                        try:
                            text = shape.TextFrame.TextRange.Text
                            has_text = bool(text and text.strip())
                        except:
                            pass
                except Exception:
                    continue

                if has_text or (text and text.strip()):
                    shape_name = "Unnamed Shape"
                    try:
                        shape_name = shape.Name
                    except:
                        pass

                    text_content[shape_id] = {"shape_name": shape_name, "text": text}
            except Exception:
                continue

        return {
            "slide_id": slide_id,
            "slide_index": slide_id,
            "slide_count": slide_count,
            "shape_count": shape_count,
            "content": text_content,
        }

    def update_text(
        self, presentation_id: str, slide_id: str, shape_id: str, text: str
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]

        try:
            clean_slide_id = slide_id.strip("\"'`")
            clean_shape_id = shape_id.strip("\"'`")
            slide_idx = int(clean_slide_id)
            shape_idx = int(clean_shape_id)
        except ValueError as e:
            return {"error": f"Invalid ID format: {str(e)}"}

        if slide_idx < 1 or slide_idx > pres.Slides.Count:
            return {"error": f"Invalid slide ID: {slide_id}"}

        try:
            slide = pres.Slides.Item(slide_idx)
        except Exception as e:
            return {"error": f"Error accessing slide: {str(e)}"}

        if shape_idx < 1 or shape_idx > slide.Shapes.Count:
            return {"error": f"Invalid shape ID: {shape_id}"}

        try:
            shape = slide.Shapes.Item(shape_idx)

            if hasattr(shape, "TextFrame2") and shape.TextFrame2.HasText:
                shape.TextFrame2.TextRange.Text = text
                return {"success": True, "message": "Text updated successfully"}
            elif hasattr(shape, "TextFrame") and hasattr(shape.TextFrame, "TextRange"):
                shape.TextFrame.TextRange.Text = text
                return {"success": True, "message": "Text updated successfully"}
            else:
                return {
                    "success": False,
                    "message": "Shape does not contain editable text",
                }
        except Exception as e:
            return {"success": False, "error": f"Error updating text: {str(e)}"}

    def add_text_box(
        self,
        presentation_id: str,
        slide_id: str,
        text: str,
        left: float = 100,
        top: float = 100,
        width: float = 400,
        height: float = 200,
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]

        try:
            clean_slide_id = slide_id.strip("\"'`")
            slide_idx = int(clean_slide_id)
        except ValueError as e:
            return {"error": f"Invalid slide ID format: {str(e)}"}

        if slide_idx < 1 or slide_idx > pres.Slides.Count:
            return {"error": f"Invalid slide ID: {slide_id}"}

        slide = pres.Slides.Item(slide_idx)

        try:
            shape = slide.Shapes.AddTextbox(1, left, top, width, height)
            shape.TextFrame.TextRange.Text = text

            shape_id = None
            for i in range(1, slide.Shapes.Count + 1):
                if slide.Shapes.Item(i) == shape:
                    shape_id = str(i)
                    break

            return {
                "success": True,
                "slide_id": slide_id,
                "shape_id": shape_id,
                "message": "Text box added successfully",
            }
        except Exception as e:
            return {"error": f"Error adding text box: {str(e)}"}

    def set_slide_title(
        self, presentation_id: str, slide_id: str, title: str
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]

        try:
            slide_idx = int(slide_id.strip("\"'"))

            if slide_idx < 1 or slide_idx > pres.Slides.Count:
                return {"error": f"Invalid slide ID: {slide_id}"}

            slide = pres.Slides.Item(slide_idx)

            title_found = False
            for shape in slide.Shapes:
                if shape.Type == 14:  # msoPlaceholder
                    if (
                        hasattr(shape, "PlaceholderFormat")
                        and shape.PlaceholderFormat.Type == 1
                    ):  # ppPlaceholderTitle
                        if hasattr(shape, "TextFrame") and hasattr(
                            shape.TextFrame, "TextRange"
                        ):
                            shape.TextFrame.TextRange.Text = title
                            title_found = True
                            break

            if not title_found:
                shape = slide.Shapes.AddTextbox(1, 50, 50, 600, 50)
                shape.TextFrame.TextRange.Text = title
                shape.TextFrame.TextRange.Font.Size = 44
                shape.TextFrame.TextRange.Font.Bold = True

            return {"success": True, "message": "Slide title has been set"}
        except Exception as e:
            return {"error": f"Error setting slide title: {str(e)}"}


class CrossPlatformPPTXAdapter(PowerPointAdapter):
    """Cross-platform adapter using python-pptx library for file-based operations"""

    def __init__(self):
        self.presentations = {}
        self.next_slide_id = 1

        # Import python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.PP_ALIGN = PP_ALIGN
            self.MSO_SHAPE_TYPE = MSO_SHAPE_TYPE
            self.available = True
        except ImportError:
            self.available = False

    def initialize(self) -> bool:
        return self.available

    def get_open_presentations(self) -> List[Dict[str, Any]]:
        result = []
        for pres_id, pres_data in self.presentations.items():
            result.append(
                {
                    "id": pres_id,
                    "name": pres_data["name"],
                    "path": pres_data["path"],
                    "slide_count": len(pres_data["presentation"].slides),
                }
            )
        return result

    def open_presentation(self, path: str) -> Dict[str, Any]:
        if not os.path.exists(path):
            return {"error": f"File not found: {path}"}

        try:
            pres = self.Presentation(path)
            pres_id = str(uuid.uuid4())
            self.presentations[pres_id] = {
                "presentation": pres,
                "name": os.path.basename(path),
                "path": path,
                "modified": False,
            }

            return {
                "id": pres_id,
                "name": os.path.basename(path),
                "path": path,
                "slide_count": len(pres.slides),
            }
        except Exception as e:
            return {"error": str(e)}

    def create_presentation(self) -> Dict[str, Any]:
        try:
            pres = self.Presentation()
            pres_id = str(uuid.uuid4())
            self.presentations[pres_id] = {
                "presentation": pres,
                "name": "New Presentation",
                "path": "",
                "modified": True,
            }

            return {
                "id": pres_id,
                "name": "New Presentation",
                "path": "",
                "slide_count": len(pres.slides),
            }
        except Exception as e:
            return {"error": str(e)}

    def save_presentation(
        self, presentation_id: str, path: str = None
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres_data = self.presentations[presentation_id]
        pres = pres_data["presentation"]

        try:
            save_path = path if path else pres_data["path"]
            if not save_path:
                return {"error": "No path specified and no original path available"}

            pres.save(save_path)
            pres_data["path"] = save_path
            pres_data["modified"] = False

            return {"success": True, "path": save_path}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def close_presentation(
        self, presentation_id: str, save: bool = True
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        if save and self.presentations[presentation_id]["modified"]:
            result = self.save_presentation(presentation_id)
            if not result.get("success", False):
                return {"success": False, "error": "Failed to save before closing"}

        del self.presentations[presentation_id]
        return {"success": True}

    def get_slides(self, presentation_id: str) -> List[Dict[str, Any]]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres = self.presentations[presentation_id]["presentation"]
        slides = []

        try:
            for i, slide in enumerate(pres.slides, 1):
                title = self._get_slide_title_pptx(slide)
                slides.append(
                    {
                        "id": str(i),
                        "index": i,
                        "title": title,
                        "shape_count": len(slide.shapes),
                    }
                )

            return slides
        except Exception as e:
            return {"error": f"Error getting slides: {str(e)}"}

    def _get_slide_title_pptx(self, slide):
        """Extract slide title from python-pptx slide object"""
        try:
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title"):
                if slide.shapes.title and hasattr(slide.shapes.title, "text"):
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        return title_text

            # Fallback: look for any text in the first shape
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    return shape.text.strip()
        except:
            pass

        return "Untitled Slide"

    def add_slide(self, presentation_id: str, layout_type: int = 1) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres_data = self.presentations[presentation_id]
        pres = pres_data["presentation"]

        try:
            # Map layout_type to python-pptx layout (simplified mapping)
            layout_map = {
                1: 0,  # Title slide
                2: 1,  # Title and content
                3: 3,  # Two content
                7: 6,  # Blank
            }
            pptx_layout = layout_map.get(layout_type, 1)

            slide_layout = pres.slide_layouts[pptx_layout]
            slide = pres.slides.add_slide(slide_layout)
            slide_index = len(pres.slides)

            pres_data["modified"] = True

            return {
                "id": str(slide_index),
                "index": slide_index,
                "title": "New Slide",
                "shape_count": len(slide.shapes),
            }
        except Exception as e:
            return {"error": f"Error adding slide: {str(e)}"}

    def get_slide_text(self, presentation_id: str, slide_id: int) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": f"Presentation ID not found: {presentation_id}"}

        pres = self.presentations[presentation_id]["presentation"]
        slide_count = len(pres.slides)

        if slide_count == 0:
            return {"error": "Presentation has no slides"}

        if slide_id < 1 or slide_id > slide_count:
            return {
                "error": f"Invalid slide ID: {slide_id}. Valid range is 1-{slide_count}"
            }

        try:
            slide = pres.slides[slide_id - 1]  # python-pptx uses 0-based indexing
            text_content = {}

            for shape_idx, shape in enumerate(slide.shapes, 1):
                if hasattr(shape, "text") and shape.text.strip():
                    shape_name = getattr(shape, "name", f"Shape {shape_idx}")
                    text_content[str(shape_idx)] = {
                        "shape_name": shape_name,
                        "text": shape.text,
                    }

            return {
                "slide_id": slide_id,
                "slide_index": slide_id,
                "slide_count": slide_count,
                "shape_count": len(slide.shapes),
                "content": text_content,
            }
        except Exception as e:
            return {
                "error": f"An error occurred: {str(e)}",
                "presentation_id": presentation_id,
                "slide_id": slide_id,
            }

    def update_text(
        self, presentation_id: str, slide_id: str, shape_id: str, text: str
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres_data = self.presentations[presentation_id]
        pres = pres_data["presentation"]

        try:
            clean_slide_id = slide_id.strip("\"'`")
            clean_shape_id = shape_id.strip("\"'`")
            slide_idx = int(clean_slide_id) - 1  # python-pptx uses 0-based indexing
            shape_idx = int(clean_shape_id) - 1
        except ValueError as e:
            return {"error": f"Invalid ID format: {str(e)}"}

        if slide_idx < 0 or slide_idx >= len(pres.slides):
            return {"error": f"Invalid slide ID: {slide_id}"}

        try:
            slide = pres.slides[slide_idx]
        except Exception as e:
            return {"error": f"Error accessing slide: {str(e)}"}

        if shape_idx < 0 or shape_idx >= len(slide.shapes):
            return {"error": f"Invalid shape ID: {shape_id}"}

        try:
            shape = slide.shapes[shape_idx]

            if hasattr(shape, "text"):
                shape.text = text
                pres_data["modified"] = True
                return {"success": True, "message": "Text updated successfully"}
            else:
                return {
                    "success": False,
                    "message": "Shape does not contain editable text",
                }
        except Exception as e:
            return {"success": False, "error": f"Error updating text: {str(e)}"}

    def add_text_box(
        self,
        presentation_id: str,
        slide_id: str,
        text: str,
        left: float = 100,
        top: float = 100,
        width: float = 400,
        height: float = 200,
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres_data = self.presentations[presentation_id]
        pres = pres_data["presentation"]

        try:
            clean_slide_id = slide_id.strip("\"'`")
            slide_idx = int(clean_slide_id) - 1  # python-pptx uses 0-based indexing
        except ValueError as e:
            return {"error": f"Invalid slide ID format: {str(e)}"}

        if slide_idx < 0 or slide_idx >= len(pres.slides):
            return {"error": f"Invalid slide ID: {slide_id}"}

        try:
            slide = pres.slides[slide_idx]

            # Convert points to inches (PowerPoint uses points, python-pptx uses inches)
            left_inches = self.Inches(left / 72)
            top_inches = self.Inches(top / 72)
            width_inches = self.Inches(width / 72)
            height_inches = self.Inches(height / 72)

            textbox = slide.shapes.add_textbox(
                left_inches, top_inches, width_inches, height_inches
            )
            textbox.text = text

            pres_data["modified"] = True
            shape_id = str(len(slide.shapes))

            return {
                "success": True,
                "slide_id": slide_id,
                "shape_id": shape_id,
                "message": "Text box added successfully",
            }
        except Exception as e:
            return {"error": f"Error adding text box: {str(e)}"}

    def set_slide_title(
        self, presentation_id: str, slide_id: str, title: str
    ) -> Dict[str, Any]:
        if presentation_id not in self.presentations:
            return {"error": "Presentation ID not found"}

        pres_data = self.presentations[presentation_id]
        pres = pres_data["presentation"]

        try:
            slide_idx = (
                int(slide_id.strip("\"'")) - 1
            )  # python-pptx uses 0-based indexing

            if slide_idx < 0 or slide_idx >= len(pres.slides):
                return {"error": f"Invalid slide ID: {slide_id}"}

            slide = pres.slides[slide_idx]

            # Try to use the slide's title placeholder
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            else:
                # Add a text box as title if no title placeholder exists
                title_textbox = slide.shapes.add_textbox(
                    self.Inches(0.5), self.Inches(0.5), self.Inches(8), self.Inches(1)
                )
                title_textbox.text = title

                # Format as title
                if hasattr(title_textbox.text_frame, "paragraphs"):
                    paragraph = title_textbox.text_frame.paragraphs[0]
                    if hasattr(paragraph, "font"):
                        paragraph.font.size = self.Pt(44)
                        paragraph.font.bold = True

            pres_data["modified"] = True

            return {"success": True, "message": "Slide title has been set"}
        except Exception as e:
            return {"error": f"Error setting slide title: {str(e)}"}


def get_powerpoint_adapter() -> PowerPointAdapter:
    """Factory function to get the appropriate PowerPoint adapter for the current platform"""
    current_platform = platform.system().lower()

    # On Windows, prefer COM adapter if available, fallback to python-pptx
    if current_platform == "windows":
        com_adapter = WindowsCOMAdapter()
        if com_adapter.available and com_adapter.initialize():
            return com_adapter
        else:
            pptx_adapter = CrossPlatformPPTXAdapter()
            if pptx_adapter.available:
                return pptx_adapter

    # On macOS and Linux, use python-pptx adapter
    pptx_adapter = CrossPlatformPPTXAdapter()
    if pptx_adapter.available:
        return pptx_adapter

    # If no adapter is available, return a placeholder
    class UnavailableAdapter(PowerPointAdapter):
        def __init__(self):
            self.error_msg = "No PowerPoint adapter available. Please install python-pptx or use Windows with PowerPoint."

        def _error_response(self):
            return {"error": self.error_msg}

        def initialize(self):
            return False

        def get_open_presentations(self):
            return self._error_response()

        def open_presentation(self, path):
            return self._error_response()

        def create_presentation(self):
            return self._error_response()

        def save_presentation(self, presentation_id, path=None):
            return self._error_response()

        def close_presentation(self, presentation_id, save=True):
            return self._error_response()

        def get_slides(self, presentation_id):
            return self._error_response()

        def add_slide(self, presentation_id, layout_type=1):
            return self._error_response()

        def get_slide_text(self, presentation_id, slide_id):
            return self._error_response()

        def update_text(self, presentation_id, slide_id, shape_id, text):
            return self._error_response()

        def add_text_box(
            self,
            presentation_id,
            slide_id,
            text,
            left=100,
            top=100,
            width=400,
            height=200,
        ):
            return self._error_response()

        def set_slide_title(self, presentation_id, slide_id, title):
            return self._error_response()

    return UnavailableAdapter()
